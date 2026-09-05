import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path="build/IntentLock_Pitch_Deck.pptx"):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Executive Razorpay & High-Tech Dark Palette
    TEXT_MAIN = RGBColor(248, 250, 252)       # Crisp Pure White (#F8FAFC)
    TEXT_MUTED = RGBColor(156, 175, 203)     # Clean Slate (#9CB0CB)
    TEXT_DIM = RGBColor(100, 116, 139)       # Dim Slate (#64748B)
    BLUE_ACCENT = RGBColor(56, 140, 255)     # Royal Blue (#388CFF)
    CYAN_ACCENT = RGBColor(0, 210, 255)      # Electric Cyan (#00D2FF)
    GREEN_ACC = RGBColor(16, 217, 144)       # Emerald Green (#10D990)
    AMBER_ACC = RGBColor(255, 179, 0)        # Amber Warning (#FFB300)
    RED_ACC = RGBColor(255, 87, 87)          # Red Alert (#FF5757)
    CARD_BG = RGBColor(13, 22, 40)           # Deep Sleek Navy (#0D1628)
    CARD_BORDER = RGBColor(30, 48, 80)       # Refined 1px Border (#1E3050)
    DIVIDER_COLOR = RGBColor(25, 40, 68)     # Header separator (#192844)

    bg_hero_path = "build/assets/slide_bg_hero.png"
    bg_content_path = "build/assets/slide_bg_content.png"

    def set_slide_background(slide, is_hero=False):
        bg_file = bg_hero_path if is_hero else bg_content_path
        if os.path.exists(bg_file):
            slide.shapes.add_picture(bg_file, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        else:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(7, 12, 22)
            bg.line.fill.background()

    def add_header(slide, title_text, category_text="RAZORPAY AI BUILDATHON 2026 • TRACK 2: AI RISK MANAGER"):
        # Category pill badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(11.7), Inches(0.32))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.size = Pt(10.5)
        p_c.font.bold = True
        p_c.font.color.rgb = CYAN_ACCENT

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.7), Inches(0.75))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN

        # Subtle clean separator bar (safe below text frame)
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.58), Inches(11.733), Inches(0.015))
        sep.fill.solid()
        sep.fill.fore_color.rgb = DIVIDER_COLOR
        sep.line.fill.background()

    def add_card(slide, left, top, width, height, title, body_lines, accent_color=BLUE_ACCENT):
        """Clean modern card container without noisy bracket tags."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1.0)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.26)
        tf.margin_right = Inches(0.26)
        tf.margin_top = Inches(0.24)
        tf.margin_bottom = Inches(0.22)

        # Card Title
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = accent_color
        p0.space_after = Pt(10)

        # Card Bullets
        for line in body_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11.5)
            p.font.color.rgb = TEXT_MAIN if not line.startswith("  ") else TEXT_MUTED
            p.space_after = Pt(4.5)

        return shape

    def add_fitted_image(slide, img_path, left, top, max_w, max_h):
        """Fit image cleanly within bounding box with proper padding and no overflow."""
        if not os.path.exists(img_path):
            return None

        # Container card behind the image
        container = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, max_w, max_h)
        container.fill.solid()
        container.fill.fore_color.rgb = CARD_BG
        container.line.color.rgb = CARD_BORDER
        container.line.width = Pt(1.0)

        with Image.open(img_path) as im:
            orig_w, orig_h = im.size

        pad_x = Inches(0.20)
        pad_y = Inches(0.20)
        inner_w = max_w - (pad_x * 2)
        inner_h = max_h - (pad_y * 2)

        img_aspect = orig_w / orig_h
        inner_aspect = inner_w / inner_h

        if img_aspect > inner_aspect:
            final_w = inner_w
            final_h = inner_w / img_aspect
        else:
            final_h = inner_h
            final_w = inner_h * img_aspect

        final_left = left + pad_x + (inner_w - final_w) / 2
        final_top = top + pad_y + (inner_h - final_h) / 2

        return slide.shapes.add_picture(img_path, final_left, final_top, width=final_w, height=final_h)

    # =========================================================================
    # SLIDE 1: Hero / Title Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, is_hero=True)

    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.35), Inches(11.333), Inches(4.8))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "RAZORPAY AI BUILDATHON 2026  •  TRACK 2: AI RISK MANAGER"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = CYAN_ACCENT
    p_badge.space_after = Pt(16)

    p_title = tf1.add_paragraph()
    p_title.text = "IntentLock"
    p_title.font.size = Pt(58)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_MAIN
    p_title.space_after = Pt(8)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Zero-Trust Policy Gate for High-Risk Payouts"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = BLUE_ACCENT
    p_sub.space_after = Pt(18)

    p_desc = tf1.add_paragraph()
    p_desc.text = "A defense-grade policy gate intercepting deepfake voice notes, urgent social engineering, and unauthorized recipient tampering before instructions reach RazorpayX rails."
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_after = Pt(28)

    # Key metric badges on hero slide
    p_pills = tf1.add_paragraph()
    p_pills.text = "403 Passing Tests  •  0 Unsafe Handoffs  •  3 Hard Authority Lanes  •  MIT Open Source"
    p_pills.font.size = Pt(13)
    p_pills.font.bold = True
    p_pills.font.color.rgb = GREEN_ACC
    p_pills.space_after = Pt(20)

    p_meta = tf1.add_paragraph()
    p_meta.text = "Builder: Kunal Kumar Mehta   •   Repository: github.com/KunalKumarMehta/razorpay"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = TEXT_DIM

    # =========================================================================
    # SLIDE 2: The Urgent Problem Landscape
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_header(s2, "The Threat Landscape: Deepfakes, Urgency & Autonomous LLM Risks")

    card_y = Inches(1.80)
    card_h = Inches(5.15)

    add_card(s2, Inches(0.8), card_y, Inches(3.64), card_h,
             "1. Weaponized Voice Clones",
             [
                 "• Rapid Voice Synthesis:",
                 "  Executive voices cloned in 30 seconds via generative audio tools.",
                 "",
                 "• Urgent WhatsApp Voice Notes:",
                 "  'Clear ₹4.25L tooling deposit to vendor account before 5 PM.'",
                 "",
                 "• Psychological Vulnerability:",
                 "  Urgency induces operators to bypass slow bank registry verification.",
                 "",
                 "• Irreversible Rails:",
                 "  Indian rails (IMPS/UPI) have zero recall: no Ctrl+Z once dispatched."
             ], RED_ACC)

    add_card(s2, Inches(4.84), card_y, Inches(3.64), card_h,
             "2. The Autonomous Agent Trap",
             [
                 "• The Industry Hype:",
                 "  Startups giving LLMs direct API keys to move company treasury.",
                 "",
                 "• The Core Fallacy:",
                 "  LLMs hallucinate, suffer prompt injection, and misunderstand edge cases.",
                 "",
                 "• Catastrophic Financial Impact:",
                 "  In enterprise payments, a 1% failure rate is an existential disaster.",
                 "",
                 "• IntentLock Thesis:",
                 "  Never let an LLM touch the financial trigger."
             ], AMBER_ACC)

    add_card(s2, Inches(8.88), card_y, Inches(3.64), card_h,
             "3. The Maker-Checker Blind Spot",
             [
                 "• Existing Protection:",
                 "  RazorpayX provides robust internal Maker-Checker queues.",
                 "",
                 "• The Upstream Hole:",
                 "  Spoofed instructions enter the queue as 'Maker Drafts' without origin validation.",
                 "",
                 "• Fatigued Reviewers:",
                 "  Approvers click 'Confirm' assuming the desk verified the counterparty.",
                 "",
                 "• The Solution:",
                 "  A zero-trust cryptographic gate BEFORE instructions enter RazorpayX."
             ], BLUE_ACCENT)

    # =========================================================================
    # SLIDE 3: Architecture & Authority Lanes (With 4K Vector Diagram)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_header(s3, "System Architecture: 3 Hard-Bounded Authority Lanes")

    arch_img = "build/diagrams/IntentLock_Architecture_Dark.png"
    if not os.path.exists(arch_img):
        arch_img = "build/diagrams/IntentLock_Architecture.png"

    # Diagram on left side
    add_fitted_image(s3, arch_img, Inches(0.8), card_y, Inches(6.8), card_h)

    # Descriptive card on right side
    add_card(s3, Inches(7.85), card_y, Inches(4.68), card_h,
             "Authority Separation Architecture",
             [
                 "• Lane 1: Trust Agent (Read-Only LLM)",
                 "  Ingests audio/images under DPDP admission rules. Extracts intent & links exact evidence spans. Zero money authority.",
                 "",
                 "• Lane 2: Deterministic Policy Gate",
                 "  Operator freezes SHA-256 intent hash. Deterministic Python core evaluates limits & approved counterparty registry.",
                 "",
                 "• Lane 3: Maker-Checker Rail",
                 "  Issues single-use 15-minute HMAC-SHA256 grant. Dispatches exactly one pending item to RazorpayX queue.",
                 "",
                 "• Cryptographic Integrity Invariant:",
                 "  Any edit to amount or account breaks the intent hash and instantly cancels the handoff grant."
             ], CYAN_ACCENT)

    # =========================================================================
    # SLIDE 4: Interactive Workflow & Fault Tolerance (With Sequence Diagram)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_header(s4, "Operational Flow: 3 Human Gestures & Fault-Tolerant Replay Prevention")

    seq_img = "build/diagrams/IntentLock_SequenceFlow_Dark.png"
    if not os.path.exists(seq_img):
        seq_img = "build/diagrams/IntentLock_SequenceFlow.png"

    # Sequence diagram on left side
    add_fitted_image(s4, seq_img, Inches(0.8), card_y, Inches(6.8), card_h)

    # Descriptive card on right side
    add_card(s4, Inches(7.85), card_y, Inches(4.68), card_h,
             "Operational Gating & Fault Tolerance",
             [
                 "• Gesture 1: Freeze Intent",
                 "  Operator verifies extracted spans. Freezes immutable SHA-256 Intent Hash into audit ledger.",
                 "",
                 "• Policy Step-Up: STEP_UP_REQUIRED",
                 "  Unapproved destination! Gate halts execution; requires independent callback + separate controller sign-off.",
                 "",
                 "• Gestures 2 & 3: Out-of-Band & Approval",
                 "  Controller approves destination -> Case transitions to ELIGIBLE. Operator freshly initiates handoff to RazorpayX.",
                 "",
                 "• Uncertainty Never Becomes Permission:",
                 "  Gateway timeouts enter RECONCILIATION_REQUIRED. Single-use HMAC burned; blind retries rejected."
             ], BLUE_ACCENT)

    # =========================================================================
    # SLIDE 5: Cryptographic Audit Ledger & Governance
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_header(s5, "Zero-Trust Cryptographic Ledger & Data Privacy Gates")

    add_card(s5, Inches(0.8), card_y, Inches(3.64), card_h,
             "Authoritative Hash Chain",
             [
                 "• Dedicated Audit Store:",
                 "  audit_events table is the sole authoritative audit store; chained via SHA-256.",
                 "",
                 "• Tip Integrity Verification:",
                 "  Verified tip checkpoints signed with HMAC-SHA256 secret keys.",
                 "",
                 "• Tamper Detection:",
                 "  Any row modification, deletion, or truncation causes immediate AuditLedgerIntegrityError.",
                 "",
                 "• Zero Mutable Leakage:",
                 "  State mutations never bypass cryptographic chain."
             ], CYAN_ACCENT)

    add_card(s5, Inches(4.84), card_y, Inches(3.64), card_h,
             "DPDP Privacy & Admission",
             [
                 "• Processing Authority Records:",
                 "  Mandatory admission record required before risk case creation.",
                 "",
                 "• Data Classification:",
                 "  Validates classification, submitter role, and retention schedules.",
                 "",
                 "• Admission Rejection:",
                 "  Invalid or unauthorized evidence triggers immediate Admission Rejection.",
                 "",
                 "• Defense-Only Architecture:",
                 "  Strictly zero offensive tooling; strictly zero PII leaks to disk."
             ], GREEN_ACC)

    add_card(s5, Inches(8.88), card_y, Inches(3.64), card_h,
             "Single-Use HMAC Grants",
             [
                 "• Single-Use HMAC Tokens:",
                 "  HMAC-SHA256 handoff grants valid for exactly 15 minutes.",
                 "",
                 "• Cryptographic Pinning:",
                 "  Bound strictly to case_id, case_version, and exact frozen intent_hash.",
                 "",
                 "• Atomic Claiming:",
                 "  Executed with SQLite BEGIN IMMEDIATE transactions to prevent race conditions.",
                 "",
                 "• Anti-Replay Guarantee:",
                 "  Tokens cannot be reused or replayed across sessions."
             ], BLUE_ACCENT)

    # =========================================================================
    # SLIDE 6: Honest Evaluation Rigor
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_header(s6, "Engineering Rigor: 403 Passing Tests & Disclosed Invariants")

    add_card(s6, Inches(0.8), card_y, Inches(5.67), card_h,
             "Evaluation Harness & Safety Suites",
             [
                 "• 45-Case Development Policy Harness:",
                 "  Validates span extraction, intent freezing, and basic registry checks.",
                 "",
                 "• 90-Case Sealed Policy Plumbing Suite:",
                 "  Covers edge cases: multi-currency, unapproved vendors, split amounts, ambiguous notes, and timeout forks.",
                 "",
                 "• 81-Execution Critical Safety Harness (27 base x 3 repetitions):",
                 "  Rigorous verification of zero Unsafe Handoffs, replay rejections, and grant expiration.",
                 "",
                 "• 403 Automated Unit & Integration Tests passing 100% in GitHub Actions CI."
             ], BLUE_ACCENT)

    add_card(s6, Inches(6.87), card_y, Inches(5.66), card_h,
             "Honest Engineering Disclosure",
             [
                 "• The Golden Rule of Hackathons: Never fake production metrics.",
                 "",
                 "• Explicit Benchmarking Declaration:",
                 "  Current benchmarks verify deterministic policy plumbing and safety boundary invariants against synthetic datasets.",
                 "",
                 "• Held-Out Production Targets (Declared as Targets):",
                 "  - Unsafe Handoffs: 0 (Zero Tolerance - Verified in Invariant Suite)",
                 "  - 3-Action Correctness Target: >= 90.0%",
                 "  - Protective Intervention Recall Target: >= 95.0%",
                 "  - Operator Interaction Reduction Target: >= 30.0%",
                 "",
                 "• Live Whisper / ASR models staged for Design Partner phase."
             ], GREEN_ACC)

    # =========================================================================
    # SLIDE 7: Startup Wedge & Business Model
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_header(s7, "Startup Thesis: Defensible Wedge for RazorpayX Payouts")

    add_card(s7, Inches(0.8), card_y, Inches(3.64), card_h,
             "1. Beachhead & Ideal Buyer",
             [
                 "• Target Buyer:",
                 "  Finance Control Owners, CFOs, & Heads of Treasury.",
                 "",
                 "• Target Segment:",
                 "  Indian growth & mid-market enterprises running on RazorpayX.",
                 "",
                 "• Urgent Wedge:",
                 "  Gating high-urgency, out-of-band payout exceptions.",
                 "",
                 "• Solves the #1 Threat:",
                 "  CEO deepfake audio inducing unauthorized wire transfers."
             ], CYAN_ACCENT)

    add_card(s7, Inches(4.84), card_y, Inches(3.64), card_h,
             "2. 90-Day Design Partner Pilot",
             [
                 "• Staged 3-Partner Pilot:",
                 "  Validate in live shadowing with 3 Indian mid-market businesses.",
                 "",
                 "• Measured Success Criteria:",
                 "  - 0 Unsafe Handoffs during live shadowing.",
                 "  - >= 30% reduction in manual operator triage time.",
                 "  - 100% catch rate on modified destination bank details.",
                 "",
                 "• Disciplined Expansion:",
                 "  Only expand into vendor onboarding if 2 of 3 partners commit."
             ], BLUE_ACCENT)

    add_card(s7, Inches(8.88), card_y, Inches(3.64), card_h,
             "3. Defensible Moat",
             [
                 "• Models are Commodities:",
                 "  Foundation models (Whisper, Gemini, Llama) get swapped constantly.",
                 "",
                 "• The Durable Moat:",
                 "  - Counterparty trust graph & approved beneficiary registry.",
                 "  - Tamper-evident cryptographic audit ledger.",
                 "  - Deep workflow embedding into RazorpayX Maker-Checker rails.",
                 "  - Provable safety invariant compliance."
             ], GREEN_ACC)

    # =========================================================================
    # SLIDE 8: Summary & The Ask
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_header(s8, "Why IntentLock Wins Track 2: AI Risk Manager", "SUMMARY & NEXT STEPS")

    add_card(s8, Inches(0.8), card_y, Inches(3.64), card_h,
             "Production-Grade Delivery",
             [
                 "• 403 passing unit & integration tests.",
                 "• GitHub Actions CI 100% Green (lint, tests, evals, Docker).",
                 "• CycloneDX & SPDX SBOM generator.",
                 "• Pinned multi-stage Docker build.",
                 "• Full FastAPI backend + Vite React Operator Console.",
                 "• MIT Open Source License."
             ], GREEN_ACC)

    add_card(s8, Inches(4.84), card_y, Inches(3.64), card_h,
             "Strategic Razorpay Alignment",
             [
                 "• Native extension for RazorpayX Payouts & Vendor Payments.",
                 "• Protects Razorpay merchants from modern deepfake threats.",
                 "• Embeds seamlessly into existing Maker-Checker approval flows.",
                 "• Strictly defense-only: 100% compliant with Track 2 rules."
             ], CYAN_ACCENT)

    add_card(s8, Inches(8.88), card_y, Inches(3.64), card_h,
             "The Ask to Judges",
             [
                 "• Architectural Grilling:",
                 "  Evaluate IntentLock on its disclosed safety invariants, code architecture, and honest metrics.",
                 "",
                 "• Razorpay Ecosystem Integration:",
                 "  Introductions to 3 RazorpayX design partners to validate in live shadowing.",
                 "",
                 "Repository:",
                 "github.com/KunalKumarMehta/razorpay"
             ], BLUE_ACCENT)

    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
